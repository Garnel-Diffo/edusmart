import io

import httpx
from docx import Document as DocxDocument
from pptx import Presentation
from pypdf import PdfReader

from app.services import groq_client


async def download_document(url: str) -> bytes:
    async with httpx.AsyncClient(timeout=30.0) as client:
        response = await client.get(url)
        response.raise_for_status()
        return response.content


def extract_text_pdf(data: bytes) -> str:
    reader = PdfReader(io.BytesIO(data))
    return "\n\n".join(page.extract_text() or "" for page in reader.pages)


def extract_text_pptx(data: bytes) -> str:
    presentation = Presentation(io.BytesIO(data))
    textes = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                textes.append("\n".join(p.text for p in shape.text_frame.paragraphs))
    return "\n\n".join(t for t in textes if t.strip())


def extract_text_docx(data: bytes) -> str:
    document = DocxDocument(io.BytesIO(data))
    return "\n\n".join(p.text for p in document.paragraphs if p.text.strip())


IMAGE_MIME_PAR_EXTENSION = {".png": "image/png", ".jpeg": "image/jpeg", ".jpg": "image/jpeg"}


async def extract_text_image(data: bytes, content_type: str = "image/jpeg") -> str:
    """OCR via LLM multimodal (UC14 étendu) : encode l'image en data-URI base64 et la transcrit."""
    import base64

    data_uri = f"data:{content_type};base64,{base64.b64encode(data).decode('ascii')}"
    return await groq_client.transcrire_image(data_uri)


EXTRACTORS = {
    "PDF": extract_text_pdf,
    "PPTX": extract_text_pptx,
    "DOCX": extract_text_docx,
}


async def extract_text(data: bytes, format_: str, content_type: str = "image/jpeg") -> str:
    if format_ == "IMAGE":
        return await extract_text_image(data, content_type)
    extractor = EXTRACTORS.get(format_)
    if extractor is None:
        raise ValueError(f"Format non supporté : {format_}")
    return extractor(data)
